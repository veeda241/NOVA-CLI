from __future__ import annotations

import csv
import os
import zipfile
from pathlib import Path


class FileReader:
    async def read(self, path: str) -> dict:
        ext = Path(path).suffix.lower()
        try:
            if ext == ".pdf":
                return self._read_pdf(path)
            if ext in {".docx", ".doc"}:
                return self._read_docx(path)
            if ext in {".xlsx", ".xls"}:
                return self._read_excel(path)
            if ext in {".pptx", ".ppt"}:
                return self._read_pptx(path)
            if ext == ".csv":
                return self._read_csv(path)
            if ext in {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff"}:
                return self._read_image(path)
            if ext in {".mp3", ".wav", ".flac", ".aac", ".mp4", ".mov", ".avi", ".mkv"}:
                return self._read_media(path)
            if ext == ".zip":
                return self._preview_zip(path)
            return self._read_text(path)
        except Exception as exc:
            return {"success": False, "error": str(exc), "path": path}

    def _read_pdf(self, path: str) -> dict:
        import pdfplumber

        text = ""
        with pdfplumber.open(path) as pdf:
            pages = len(pdf.pages)
            for page in pdf.pages[:20]:
                text += (page.extract_text() or "") + "\n"
        return {"success": True, "type": "pdf", "path": path, "pages": pages, "content": text[:50000]}

    def _read_docx(self, path: str) -> dict:
        from docx import Document

        doc = Document(path)
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
        return {"success": True, "type": "docx", "path": path, "content": text[:50000]}

    def _read_excel(self, path: str) -> dict:
        import openpyxl

        workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
        chunks: list[str] = []
        for sheet_name in workbook.sheetnames[:5]:
            sheet = workbook[sheet_name]
            chunks.append(f"--- Sheet: {sheet_name} ---")
            for idx, row in enumerate(sheet.iter_rows(values_only=True)):
                if idx >= 100:
                    break
                chunks.append(" | ".join("" if value is None else str(value) for value in row))
        return {"success": True, "type": "excel", "path": path, "content": "\n".join(chunks)[:50000]}

    def _read_pptx(self, path: str) -> dict:
        from pptx import Presentation

        presentation = Presentation(path)
        chunks: list[str] = []
        for index, slide in enumerate(presentation.slides[:30], start=1):
            chunks.append(f"--- Slide {index} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
        return {"success": True, "type": "pptx", "path": path, "content": "\n".join(chunks)[:50000]}

    def _read_csv(self, path: str) -> dict:
        rows: list[str] = []
        with open(path, "r", errors="ignore", newline="", encoding="utf-8") as handle:
            reader = csv.reader(handle)
            for index, row in enumerate(reader):
                if index >= 250:
                    break
                rows.append(" | ".join(row))
        return {"success": True, "type": "csv", "path": path, "content": "\n".join(rows)}

    def _read_image(self, path: str) -> dict:
        from PIL import Image

        image = Image.open(path)
        return {
            "success": True,
            "type": "image",
            "path": path,
            "content": f"Image: {image.format} {image.size[0]}x{image.size[1]} {image.mode}",
            "width": image.size[0],
            "height": image.size[1],
            "format": image.format,
        }

    def _read_media(self, path: str) -> dict:
        return {"success": True, "type": "media", "path": path, "content": f"Media file: {os.path.basename(path)}", "size": os.path.getsize(path)}

    def _preview_zip(self, path: str) -> dict:
        with zipfile.ZipFile(path, "r") as archive:
            names = archive.namelist()
        return {"success": True, "type": "archive", "path": path, "content": "\n".join(names[:100]), "count": len(names)}

    def _read_text(self, path: str) -> dict:
        with open(path, "r", errors="ignore", encoding="utf-8") as handle:
            content = handle.read(50000)
        return {"success": True, "type": "text", "path": path, "content": content}
